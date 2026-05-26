from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
slide_layout = prs.slide_layouts[6]  # Blank layout
slide = prs.slides.add_slide(slide_layout)

# Background color (dark blue)
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0x00, 0x33, 0x66)

def add_textbox(slide, text, left, top, width, height, font_size, bold=False, color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color

W, H = prs.slide_width, prs.slide_height

# Title
add_textbox(slide, "Welcome!", Inches(1), Inches(1.2), W - Inches(2), Inches(1.2), 54, bold=True)

# Subtitle
add_textbox(slide, "Innovating Together for a Better Tomorrow", Inches(1), Inches(2.5), W - Inches(2), Inches(0.8), 24, color=RGBColor(0xAD, 0xD8, 0xE6))

# Event name
add_textbox(slide, "AI & Technology Summit 2025", Inches(1), Inches(3.5), W - Inches(2), Inches(0.7), 20, bold=True, color=RGBColor(0xFF, 0xD7, 0x00))

# Presenter name
add_textbox(slide, "Presented by: Irshad | Ericsson", Inches(1), Inches(4.4), W - Inches(2), Inches(0.6), 16, color=RGBColor(0xCC, 0xCC, 0xCC))

# Date
add_textbox(slide, "May 2026", Inches(1), Inches(5.1), W - Inches(2), Inches(0.5), 14, color=RGBColor(0xCC, 0xCC, 0xCC))

prs.save("welcome_slide_v2.pptx")
print("welcome_slide_v2.pptx created successfully!")
